from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "docs" / "presentations" / "workshop_presentation.pptx"
BACKUP = ROOT / "docs" / "presentations" / "workshop_presentation.original-backup.pptx"
PRESENTATIONS = ROOT / "docs" / "presentations"

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(27, 45, 73)
GREEN = RGBColor(0, 104, 83)
TEAL = RGBColor(0, 122, 132)
BLUE = RGBColor(21, 101, 192)
INDIGO = RGBColor(40, 53, 147)
PURPLE = RGBColor(106, 27, 154)
CRIMSON = RGBColor(173, 20, 87)
RED = RGBColor(183, 28, 28)
BROWN = RGBColor(93, 64, 55)
DARK = RGBColor(49, 59, 73)
GREY = RGBColor(246, 248, 250)
MID_GREY = RGBColor(218, 225, 230)
TEXT = RGBColor(31, 41, 55)
WHITE = RGBColor(255, 255, 255)
AMBER = RGBColor(245, 158, 11)


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def clear_deck(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_text(shape, text, size=18, bold=False, color=TEXT, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_box(slide, x, y, w, h, text, fill=GREY, line=MID_GREY, size=15, bold=False, color=TEXT, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x,
        y,
        w,
        h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(shape, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER)
    return shape


def add_title(slide, title, section=None, number=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), H)
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.fill.background()
    if section:
        sec = slide.shapes.add_textbox(Inches(0.55), Inches(0.2), Inches(3.6), Inches(0.28))
        set_text(sec, section.upper(), size=8, bold=True, color=GREEN)
    t = slide.shapes.add_textbox(Inches(0.52), Inches(0.46), Inches(11.8), Inches(0.55))
    set_text(t, title, size=24, bold=True, color=NAVY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(1.1), Inches(11.9), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = MID_GREY
    line.line.fill.background()
    if number:
        foot = slide.shapes.add_textbox(Inches(11.95), Inches(7.05), Inches(0.9), Inches(0.25))
        set_text(foot, str(number), size=8, color=RGBColor(120, 130, 140), align=PP_ALIGN.RIGHT)


def bullet_list(slide, items, x, y, w, h, size=15, gap=0.05, color=TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.04)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap * 72)
    return box


def add_table(slide, rows, x, y, w, h, col_widths=None, font_size=10, header_fill=NAVY):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = table_shape.table
    if col_widths:
        for idx, cw in enumerate(col_widths):
            table.columns[idx].width = cw
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            fill = header_fill if r_idx == 0 else (GREY if r_idx % 2 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size if r_idx else font_size + 1)
                p.font.bold = r_idx == 0
                p.font.color.rgb = WHITE if r_idx == 0 else TEXT
    return table_shape


def connect(slide, x1, y1, x2, y2, color=RGBColor(110, 120, 130), width=1.6):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def section_slide(prs, title, subtitle, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.85), W, Inches(0.65)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = GREEN
    t = slide.shapes.add_textbox(Inches(0.75), Inches(2.35), Inches(11.5), Inches(0.8))
    set_text(t, title, size=34, bold=True, color=WHITE)
    s = slide.shapes.add_textbox(Inches(0.78), Inches(3.2), Inches(10.8), Inches(0.4))
    set_text(s, subtitle, size=17, color=RGBColor(214, 232, 228))
    f = slide.shapes.add_textbox(Inches(11.7), Inches(7.05), Inches(1), Inches(0.25))
    set_text(f, str(n), size=8, color=WHITE, align=PP_ALIGN.RIGHT)
    return slide


def add_full_slide_diagram(slide, image_name, number=None):
    path = PRESENTATIONS / image_name
    pic = slide.shapes.add_picture(str(path), 0, 0)
    scale = min(W / pic.width, H / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int((W - pic.width) / 2)
    pic.top = int((H - pic.height) / 2)
    if number:
        foot = slide.shapes.add_textbox(Inches(12.0), Inches(7.08), Inches(0.8), Inches(0.22))
        set_text(foot, str(number), size=8, color=RGBColor(120, 130, 140), align=PP_ALIGN.RIGHT)


def diagram_slide(prs, image_name, number):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_full_slide_diagram(slide, image_name, number)
    return slide


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.26)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = GREEN
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.88), W, Inches(0.62)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = GREEN
    t = slide.shapes.add_textbox(Inches(0.75), Inches(1.25), Inches(11.8), Inches(0.8))
    set_text(t, "Debt Collection Management System", size=34, bold=True, color=WHITE)
    st = slide.shapes.add_textbox(Inches(0.78), Inches(2.08), Inches(11.1), Inches(0.55))
    set_text(st, "Leadership presentation | Greenfield delivery baseline", size=18, color=RGBColor(218, 232, 230))
    mission = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(10.5), Inches(1.05))
    set_text(
        mission,
        "DCMS is a greenfield platform for managing DWP debt from identification through recovery, enforcement, and write-off, with compliance, auditability, and business-change control built into the product.",
        size=22,
        bold=True,
        color=WHITE,
    )
    add_box(slide, Inches(0.82), Inches(4.75), Inches(3.25), Inches(0.65), "Not COTS", GREEN, GREEN, 18, True, WHITE)
    add_box(slide, Inches(4.35), Inches(4.75), Inches(3.25), Inches(0.65), "Not Solon Tax", TEAL, TEAL, 18, True, WHITE)
    add_box(slide, Inches(7.88), Inches(4.75), Inches(3.25), Inches(0.65), "Built for DWP rules", BLUE, BLUE, 18, True, WHITE)


def slide_scope_picture(prs, n):
    return diagram_slide(prs, "diag-01-debt-lifecycle.png", n)


def slide_differentiators(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "What Makes This Different", "What we are building", n)
    cards = [
        (
            "Regulatory hardcoding",
            "Breathing space, insolvency, deceased, vulnerability, and fraud constraints cannot be configured out of effect.",
        ),
        (
            "Business-user strategy authoring",
            "Decision tables and treatment paths can change through governed authoring without an IT ticket for every rule change.",
        ),
        (
            "Audit-to-question trail",
            "A regulatory query can be answered from account history, rule versions, suppression decisions, and actor evidence inside DCMS.",
        ),
    ]
    colors = [CRIMSON, GREEN, BLUE]
    for i, (head, body) in enumerate(cards):
        x = Inches(0.75 + i * 4.15)
        add_box(slide, x, Inches(1.72), Inches(3.65), Inches(0.72), head, colors[i], colors[i], 15, True, WHITE)
        body_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.65), Inches(3.65), Inches(2.2))
        body_box.fill.solid()
        body_box.fill.fore_color.rgb = GREY
        body_box.line.color.rgb = MID_GREY
        body_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        set_text(body_box, body, size=14, color=TEXT, align=PP_ALIGN.LEFT)
    add_box(
        slide,
        Inches(1.1),
        Inches(5.65),
        Inches(11.0),
        Inches(0.65),
        "Leadership message: the design turns compliance intent into product guardrails.",
        GREY,
        MID_GREY,
        17,
        True,
        NAVY,
    )


def slide_architecture(prs, n):
    return diagram_slide(prs, "diag-02-system-architecture.png", n)


def slide_delegate(prs, n):
    return diagram_slide(prs, "diag-03-delegate-command-pattern.png", n)


def slide_adr_table(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Key Architecture Decisions", "Architecture overview", n)
    rows = [
        ["Decision", "Rationale / risk controlled"],
        ["Single monolith", "Preserves complexity budget for a 6-7 person team and one operational model."],
        ["One Flowable process per DebtAccount", "The debt lifecycle is the process; customer view is assembled across accounts."],
        ["Java 21 exhaustive switch", "Compile-time guarantee for suppression coverage where a miss could become a legal breach."],
        ["Transaction boundary rule", "DB writes commit before Flowable calls; protective state is never rolled back by engine failure."],
        ["Write-off self-approval prevention", "Flowable, service, and DB checks jointly enforce delegated-authority controls."],
        ["Keycloak RBAC via JWT", "One OAuth2/OIDC security model; React hiding is UX only, backend APIs enforce authority."],
    ]
    add_table(slide, rows, Inches(0.75), Inches(1.48), Inches(11.85), Inches(4.8), [Inches(3.1), Inches(8.75)], 12)


def slide_requirements(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Requirement Coverage Map", "Requirements coverage", n)
    rows = [
        ["Coverage band", "Capability groups", "Message"],
        ["Covered", "CAS, DIC, RPF, UAAF, MIR, A, AAD, I3PS, CP, SD", "Ownership and design materially defined."],
        ["Tier A gaps", "UI, SoR, MR", "Known design work remains: frontend catalogue, finance depth, migration tooling."],
        ["Tier B blocked", "WAM, IEC, CC, 3PM, BSF", "Design narrowed, blocked on DWP policy values or sign-off."],
        ["Tier C scaffold", "DW plus process packages", "Conceptual design exists; package/code scaffolding is sequencing work."],
    ]
    add_table(slide, rows, Inches(0.7), Inches(1.55), Inches(11.95), Inches(3.35), [Inches(2.15), Inches(4.05), Inches(5.75)], 11)
    add_box(slide, Inches(1.0), Inches(5.35), Inches(11.2), Inches(0.75), "All 19 groups have confirmed module ownership; gaps are classified, not unknown.", GREEN, GREEN, 18, True, WHITE)


def slide_gap_classification(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Gap Classification", "Requirements coverage", n)
    rows = [
        ["Tier", "Description", "Risk owner"],
        ["A", "Design not yet done: complaints, legal/enforcement, finance depth, broader frontend catalogue.", "Delivery team"],
        ["B", "Design exists but policy values/sign-off are needed: I&E staleness, DCA notice, thresholds.", "DWP programme"],
        ["C", "Scaffolding only: code/package structure exists conceptually, build not started or not wired.", "Delivery team sequencing"],
    ]
    add_table(slide, rows, Inches(1.0), Inches(1.7), Inches(11.0), Inches(2.65), [Inches(1.2), Inches(7.0), Inches(2.8)], 13, GREEN)
    bullet_list(slide, [
        "No unknown gaps: every gap is named, owned, and categorised.",
        "Older Solon/additional-layer assumptions are obsolete against the current monolith baseline.",
        "Next planning work should close targeted policy and design inputs, not reopen architecture.",
    ], Inches(1.15), Inches(4.75), Inches(10.8), Inches(1.2), 16)


def slide_implied_scope(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Implied Build Scope", "Requirements coverage", n)
    rows = [
        ["Discovered scope", "Why it is necessary"],
        ["Audit trail completeness", "Regulatory questions need answerable in-system evidence."],
        ["Configuration version snapshot", "Policy activation must not silently change live treatment paths."],
        ["Exhaustive suppression switch", "A missed statutory suppression category must fail at build time."],
        ["Three-layer write-off control", "Financial-control assurance needs workflow, service, and DB enforcement."],
        ["Deceased handler two-phase atomicity", "A Flowable failure must not leave active collection against a deceased party."],
        ["Champion/challenger simulation", "Strategy changes need test evidence before live activation."],
        ["DCA pre-placement window", "GDPR disclosure is enforced by state machine, not operator discretion."],
    ]
    add_table(slide, rows, Inches(0.7), Inches(1.35), Inches(12.0), Inches(5.45), [Inches(3.35), Inches(8.65)], 10)


def slide_entity(prs, n):
    return diagram_slide(prs, "diag-04-entity-model.png", n)


def slide_reg_state(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Regulatory State on the Account", "Customer and account model", n)
    states = [("Breathing space", CRIMSON), ("Insolvency", PURPLE), ("Deceased", DARK), ("Fraud marker", AMBER)]
    for i, (name, color) in enumerate(states):
        add_box(slide, Inches(0.85 + i * 3.05), Inches(1.85), Inches(2.45), Inches(1.05), name, color, color, 17, True, WHITE)
    bullet_list(slide, [
        "These are DB-first facts: the process engine reads them fresh at each decision point.",
        "They are not stored only as process variables, which can become stale or incomplete.",
        "No configuration can remove their legal effect from collection, communication, or audit handling.",
    ], Inches(1.1), Inches(3.65), Inches(11), Inches(1.4), 18)
    add_box(slide, Inches(1.2), Inches(5.7), Inches(10.7), Inches(0.7), "Leadership message: statutory protection is a system invariant, not a rule-table preference.", GREEN, GREEN, 17, True, WHITE)


def slide_config_tiers(prs, n):
    return diagram_slide(prs, "diag-05-three-tier-configuration-model.png", n)


def slide_policy_bundle(prs, n):
    return diagram_slide(prs, "diag-06-policy-bundle-lifecycle.png", n)


def slide_demo_matrix(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "14 Flows Covering 20 Behavioural Requirements", "Demo flow preview", n)
    flows = [
        ("1", "Intake to first contact"), ("2", "Vulnerability to resolution"), ("3", "Breach to DCA"),
        ("4", "Complex household"), ("5", "Strategy change without IT"), ("6", "Executive dashboard"),
        ("7", "Self-service handoff"), ("8", "Dispute"), ("9", "Deceased protocol"),
        ("10", "Write-off and reactivation"), ("11", "New agent journey"), ("12", "Regulatory audit"),
        ("13", "Month-end surge"), ("14", "Settlement offer"),
    ]
    left_x = Inches(0.75)
    y = Inches(1.45)
    for idx, (num, name) in enumerate(flows):
        row_y = y + idx * Inches(0.34)
        fill = GREEN if num in {"2", "5", "12"} else GREY
        color = WHITE if num in {"2", "5", "12"} else TEXT
        add_box(slide, left_x, row_y, Inches(3.3), Inches(0.25), f"Flow {num}: {name}", fill, fill if fill == GREEN else MID_GREY, 7.5, num in {"2", "5", "12"}, color, False)
        for j in range(20):
            cell_fill = rgb("D9F0E8") if (idx + j) % 4 == 0 else WHITE
            if num in {"2", "5", "12"} and j in {6, 11, 12, 17}:
                cell_fill = GREEN
            add_box(slide, Inches(4.25) + j * Inches(0.36), row_y, Inches(0.26), Inches(0.25), "x" if cell_fill != WHITE else "", cell_fill, MID_GREY, 6, True, WHITE if cell_fill == GREEN else GREEN, False)
    add_box(slide, Inches(4.25), Inches(6.55), Inches(7.25), Inches(0.45), "All 20 behavioural requirements are covered by at least one flow; headline flows are 2, 5, and 12.", NAVY, NAVY, 12, True, WHITE)


def slide_demo_moments(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Three Headline Demo Moments", "Demo flow preview", n)
    moments = [
        ("Flow 5", "Strategy change without IT", "Business admin edits DMN, creates a policy bundle, runs simulation, and activates atomically."),
        ("Flow 2", "Vulnerability to resolution", "Vulnerability flag causes suppression, specialist handling, I&E, and affordable arrangement in under 2 minutes."),
        ("Flow 12", "Regulatory audit", "Account history, rule versions, actors, and suppression decisions answer a regulatory query in under 10 minutes."),
    ]
    for i, (flow, head, body) in enumerate(moments):
        x = Inches(0.9 + i * 4.05)
        add_box(slide, x, Inches(1.75), Inches(3.45), Inches(0.6), flow, GREEN, GREEN, 17, True, WHITE)
        add_box(slide, x, Inches(2.5), Inches(3.45), Inches(0.55), head, NAVY, NAVY, 13, True, WHITE)
        b = slide.shapes.add_textbox(x + Inches(0.08), Inches(3.25), Inches(3.25), Inches(1.6))
        set_text(b, body, 14, False, TEXT)


def slide_solon_options(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Solon Tax v2.3.0: Options Reassessed", "Options, limitations, decision", n)
    rows = [
        ["Option", "Description", "Gross reuse", "Net value", "Decision"],
        ["A", "Full Solon Tax platform adoption", "30-35%", "5-10%", "Rejected: 7 hard blockers"],
        ["B", "Solon Tax financial engine only", "20-25%", "10-15%", "Rejected: blocker + integration cost"],
        ["C", "Greenfield current path", "Baseline", "Baseline", "Confirmed: zero hard blockers"],
    ]
    add_table(
        slide,
        rows,
        Inches(0.65),
        Inches(1.55),
        Inches(12.05),
        Inches(2.45),
        [Inches(0.85), Inches(4.45), Inches(1.75), Inches(1.6), Inches(3.4)],
        11,
        NAVY,
    )
    add_box(
        slide,
        Inches(0.9),
        Inches(4.35),
        Inches(5.45),
        Inches(0.85),
        "v2.3.0 resolves the old auth and database blockers, and retires the Camunda 7 EOL objection.",
        GREY,
        MID_GREY,
        13,
        True,
        NAVY,
    )
    add_box(
        slide,
        Inches(6.75),
        Inches(4.35),
        Inches(5.45),
        Inches(0.85),
        "The decision is unchanged: no Solon component is adopted as platform, process engine, financial engine, or structural dependency.",
        GREY,
        MID_GREY,
        13,
        True,
        NAVY,
    )
    add_box(
        slide,
        Inches(1.25),
        Inches(5.8),
        Inches(10.5),
        Inches(0.7),
        "Decision frame: 5-15% net scope reduction does not justify statutory, integration, or platform coupling risk.",
        CRIMSON,
        CRIMSON,
        16,
        True,
        WHITE,
    )


def slide_solon_locked(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Why the Decision Is Locked", "Options, limitations, decision", n)
    add_box(
        slide,
        Inches(0.8),
        Inches(1.35),
        Inches(11.7),
        Inches(1.15),
        "Primary disqualifier: Solon Tax's Amplio engine supports interrupting boundary events only. DCMS requires non-interrupting events for breathing space; sending debt collection communication during an active moratorium is a criminal offence.",
        CRIMSON,
        CRIMSON,
        14,
        True,
        WHITE,
    )
    blockers = [
        "No native DMN: Drools requires developer involvement.",
        "Java 17 vs Java 21: loses exhaustive switch safety.",
        "Liquibase vs Flyway: incompatible migration history.",
        "Angular vs React + GOV.UK Design System.",
        "Microservices vs intentional monolith.",
        "OPA vs Keycloak JWT RBAC.",
    ]
    bullet_list(slide, blockers, Inches(1.05), Inches(2.85), Inches(11), Inches(2.2), 14)
    add_box(
        slide,
        Inches(1.0),
        Inches(5.45),
        Inches(11.0),
        Inches(0.52),
        "Option B still carries Liquibase/Flyway risk and a long-lived anti-corruption layer for only 10-15% net value.",
        GREY,
        MID_GREY,
        13,
        True,
        NAVY,
    )
    add_box(
        slide,
        Inches(1.0),
        Inches(6.15),
        Inches(11.0),
        Inches(0.55),
        "Revisit only if DWP mandates a platform or a future Solon Tax version resolves all seven blockers.",
        GREY,
        MID_GREY,
        14,
        True,
        NAVY,
    )


def slide_release1(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Release 1 Scope", "Development releases", n)
    rows = [
        ["In scope", "Explicitly not in Release 1"],
        ["customer, debtaccount, repaymentplan, communications basic, workallocation basic", "DCA placement 9-state machine"],
        ["strategy Tier 2 DMN, audit, Keycloak RBAC, process ports/infrastructure baseline", "Legal/enforcement, analytics/reporting depth, full I&E depth"],
        ["Demo flows 1, 2, 3, 4, 6 core agent workflow", "External integrations, self-service customer API, full finance accounting depth"],
    ]
    add_table(slide, rows, Inches(0.8), Inches(1.55), Inches(11.7), Inches(3.6), [Inches(6.0), Inches(5.7)], 12, GREEN)
    add_box(slide, Inches(1.1), Inches(5.75), Inches(10.9), Inches(0.72), "Release 1 is working software in a running environment, not a feature-complete system.", NAVY, NAVY, 16, True, WHITE)


def slide_release_beyond(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "Release 2 and Beyond", "Development releases", n)
    rows = [
        ["Release", "Key additions"],
        ["Release 2", "DCA placement, full champion/challenger rollout, policy bundle UI, self-service API, finance accounting depth."],
        ["Release 3", "Legal/enforcement, complaints, analytics/reporting completion, migration tooling, test/UAT environments."],
    ]
    add_table(slide, rows, Inches(1.0), Inches(1.9), Inches(11.0), Inches(2.2), [Inches(2.0), Inches(9.0)], 14, NAVY)
    add_box(slide, Inches(1.2), Inches(5.0), Inches(10.7), Inches(0.75), "Releases are milestone gates, not waterfall phases: each gate must produce demonstrable software and evidence.", GREEN, GREEN, 16, True, WHITE)


def slide_decisions(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, "What Leadership Needs to Decide", "Close", n)
    rows = [
        ["Blocked item", "Why it blocks delivery"],
        ["I&E staleness period", "Determines when I&E must be recaptured before an arrangement."],
        ["DCA notice period", "Sets the PRE_PLACEMENT_NOTICE window in the 9-state machine."],
        ["Champion/challenger thresholds", "Defines split, duration, sample-size, and promotion gates."],
        ["Vulnerability lawful basis", "Required before vulnerability data can be processed under GDPR."],
        ["Statute-barred timing", "Determines when a debt is legally unenforceable."],
    ]
    add_table(slide, rows, Inches(0.9), Inches(1.45), Inches(11.55), Inches(3.75), [Inches(3.2), Inches(8.35)], 12, CRIMSON)
    add_box(slide, Inches(1.15), Inches(5.75), Inches(10.9), Inches(0.7), "Design is complete enough to start delivery. These five decisions unblock the corresponding features.", GREEN, GREEN, 16, True, WHITE)


def appendix_slide(prs, letter, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[31])
    add_title(slide, f"Appendix {letter}: {title}", "Feature deep dives", f"A{letter}")
    bullet_list(slide, bullets, Inches(0.95), Inches(1.55), Inches(11.5), Inches(4.8), 15)


def build():
    if not BACKUP.exists():
        shutil.copy2(PPTX, BACKUP)
    prs = Presentation(PPTX)
    clear_deck(prs)
    prs.slide_width = W
    prs.slide_height = H

    title_slide(prs)
    slide_scope_picture(prs, 2)
    slide_differentiators(prs, 3)
    slide_architecture(prs, 4)
    slide_delegate(prs, 5)
    slide_adr_table(prs, 6)
    slide_requirements(prs, 7)
    slide_gap_classification(prs, 8)
    slide_implied_scope(prs, 9)
    slide_entity(prs, 10)
    slide_reg_state(prs, 11)
    slide_config_tiers(prs, 12)
    slide_policy_bundle(prs, 13)
    slide_demo_matrix(prs, 14)
    slide_demo_moments(prs, 15)
    slide_solon_options(prs, 16)
    slide_solon_locked(prs, 17)
    slide_release1(prs, 18)
    slide_release_beyond(prs, 19)
    slide_decisions(prs, 20)

    appendix = [
        ("A", "Champion/Challenger", [
            "Two DMN versions are deployed simultaneously.",
            "Deterministic account-level routing keeps accounts on their assigned version.",
            "Comparison dashboard tracks performance; challenger promotion retires the old champion.",
            "Source: ADR-010 and strategy-simulation-engine-design.md.",
        ]),
        ("B", "Strategy Simulation Engine", [
            "Runs proposed rule changes against historical account populations before live activation.",
            "Outputs projected treatment distribution and outcomes versus current policy.",
            "Gives business users test evidence without exposing live accounts.",
        ]),
        ("C", "Policy Bundle Activation", [
            "Lifecycle: DRAFT -> READY -> APPROVED -> ACTIVE -> RETIRED.",
            "Tier 1, 2, and 3 changes share one effective date, approval, and audit entry.",
            "Any activation failure rolls back the whole bundle to APPROVED.",
        ]),
        ("D", "Vulnerability Handling", [
            "Agent-set vulnerability flag immediately suppresses outbound collection contact.",
            "Specialist routing and I&E are mandated before affordable arrangement creation.",
            "This is hardcoded compliance behaviour, not configurable discretion.",
        ]),
        ("E", "Breathing Space", [
            "Non-interrupting BPMN event subprocess applies the moratorium immediately.",
            "Collection actions suspend for the statutory period and resume on expiry.",
            "Incorrect communication during moratorium carries criminal liability.",
        ]),
        ("F", "Deceased Party Handling", [
            "Phase 1 commits deceased flag, suppression, and audit atomically.",
            "Phase 2 suspends Flowable process instances outside the transaction.",
            "A process-engine failure cannot roll back the legally protective DB state.",
        ]),
        ("G", "DCA Placement State Machine", [
            "Nine states include CANDIDATE, PRE_PLACEMENT_NOTICE, AWAITING_ACKNOWLEDGEMENT, PLACED, ACTIVE, and terminal outcomes.",
            "PRE_PLACEMENT_NOTICE enforces mandatory disclosure before placement.",
            "There is no shortcut transition directly to PLACED.",
        ]),
        ("H", "Regulatory Audit Trail", [
            "Actions, DMN invocations, suppression evaluations, process snapshots, actors, timestamps, and rule versions are logged.",
            "Single account history supports regulatory reconstruction without external manual assembly.",
            "Demo Flow 12 targets question-to-answer in under 10 minutes.",
        ]),
        ("I", "Self-Service Customer Integration", [
            "DCMS exposes backend APIs for I&E submission, payment events, arrangement requests, and contact updates.",
            "The customer portal is external to DCMS.",
            "Demo Flow 7 shows contract and workflow handoff, not a customer-facing UI.",
        ]),
    ]
    for letter, title, bullets in appendix:
        appendix_slide(prs, letter, title, bullets)

    prs.save(PPTX)
    print(f"Saved {PPTX}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
